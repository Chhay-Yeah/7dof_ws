#!/usr/bin/env python3
"""
Render the ROS 2 computation graph of the teach-pendant backend running in
GAZEBO mode, in the visual style of ``rqt_graph`` (Nodes/Topics view: nodes are
ellipses, topics are boxes, arrows are publish/subscribe).

The topology is the real one — extracted from each node's create_publisher /
create_subscription calls + the pendant_backend.launch.py graph + the standard
ros2_control / Gazebo components — minus the rqt clutter (/rosout,
/parameter_events) so it reads cleanly in a thesis.

    python3 make_system_graph.py            # -> system_rqt_graph.{png,pdf,dot}
    python3 make_system_graph.py --out figures/system_rqt_graph

Needs the Graphviz `dot` binary on PATH (no python bindings required).
"""
import argparse
import subprocess
from pathlib import Path

# ── node pub/sub (topic, direction) — direction 'pub' or 'sub' ──────────────
# Lifted from the live nodes; see pendant_backend.launch.py for the launch set.
NODES = {
    "/pendant7dof_bridge": {   # the 7dof-pendant GUI (in-process rclpy node)
        "pub": ["/ee_target", "/drawing/strokes", "/arm_controller/joint_trajectory",
                "/cartesian_path"],
        "sub": ["/joint_states", "/ee_pose", "/pen_canvas_norm"],
    },
    "/ik_7dof_v3": {           # ik_arm_v3.py — Cartesian-jog IK
        "pub": ["/joint_commands"],
        "sub": ["/ee_target", "/joint_states", "/robot_description"],
    },
    "/fk_7dof_v3": {           # fk_arm_v3.py — EE pose for status/jog
        "pub": ["/ee_pose"],
        "sub": ["/joint_states", "/robot_description"],
    },
    "/ik_to_trajectory": {     # /joint_commands -> controller
        "pub": ["/arm_controller/joint_trajectory"],
        "sub": ["/joint_commands", "/joint_states"],
    },
    "/drawing_batch_planner": {  # /drawing/strokes -> one JointTrajectory
        "pub": ["/arm_controller/joint_trajectory", "/cartesian_path",
                "/pen_canvas_norm"],
        "sub": ["/drawing/strokes", "/joint_states", "/robot_description"],
    },
    "/go_to_start": {          # boots the arm to the elbow-up start pose
        "pub": ["/arm_controller/joint_trajectory"],
        "sub": ["/joint_states"],
    },
    "/robot_state_publisher": {
        "pub": ["/robot_description", "/tf", "/tf_static"],
        "sub": ["/joint_states"],
    },
    "/joint_state_broadcaster": {"pub": ["/joint_states"], "sub": []},
    "/arm_controller": {"pub": [], "sub": ["/arm_controller/joint_trajectory"]},
    "/gz_ros2_control": {"pub": ["/clock"], "sub": []},  # Ignition + ros2_control plugin
}

# Controllers hosted inside the Gazebo controller_manager process — grouped.
GAZEBO_CLUSTER = {"/gz_ros2_control", "/joint_state_broadcaster", "/arm_controller"}

# Colours (rqt-ish): nodes light blue, topics pale yellow, sim group pale green.
C_NODE, C_NODE_LINE = "#cfe2ff", "#3a6ea5"
C_TOPIC, C_TOPIC_LINE = "#fff3c4", "#b58900"
C_GZNODE = "#cdeccd"
C_EDGE = "#5b6470"


def _esc(t):  # graphviz id
    return '"' + t + '"'


def build_dot() -> str:
    topics = sorted({t for n in NODES.values() for t in (n["pub"] + n["sub"])})
    lines = [
        "digraph ros_graph {",
        '  rankdir=LR;',
        '  bgcolor="white";',
        '  labelloc="t"; fontname="Helvetica-Bold"; fontsize=20;',
        '  label="ROS 2 computation graph — pendant backend (Gazebo mode)\\n'
        'rqt_graph Nodes/Topics view";',
        '  node [fontname="Helvetica", fontsize=12];',
        '  edge [color="%s", arrowsize=0.8, penwidth=1.1];' % C_EDGE,
        "",
        "  // ── topics (boxes) ──",
    ]
    for t in topics:
        lines.append('  %s [shape=box, style="rounded,filled", '
                     'fillcolor="%s", color="%s"];' % (_esc(t), C_TOPIC, C_TOPIC_LINE))
    lines.append("")
    lines.append("  // ── nodes (ellipses) outside the sim ──")
    for n in NODES:
        if n in GAZEBO_CLUSTER:
            continue
        lines.append('  %s [shape=ellipse, style=filled, fillcolor="%s", '
                     'color="%s"];' % (_esc(n), C_NODE, C_NODE_LINE))
    lines += [
        "",
        "  // ── Gazebo / ros2_control group ──",
        '  subgraph cluster_gz {',
        '    label="Gazebo / ros2_control  (controller_manager)";',
        '    fontname="Helvetica-Bold"; fontsize=13; color="#2e7d32";',
        '    style="rounded,filled"; fillcolor="#eef7ee";',
    ]
    for n in sorted(GAZEBO_CLUSTER):
        lines.append('    %s [shape=ellipse, style=filled, fillcolor="%s", '
                     'color="#2e7d32"];' % (_esc(n), C_GZNODE))
    lines.append("  }")
    lines.append("")
    lines.append("  // ── publish (node -> topic) & subscribe (topic -> node) ──")
    for n, io in NODES.items():
        for t in io["pub"]:
            lines.append("  %s -> %s;" % (_esc(n), _esc(t)))
        for t in io["sub"]:
            lines.append("  %s -> %s;" % (_esc(t), _esc(n)))
    lines.append("}")
    return "\n".join(lines)


def _parse_plain(dot_path):
    """Run `dot -Tplain` and return (graph_w, graph_h, nodes, edges) in inches.
    nodes: name -> (cx, cy, w, h)  [dot coords, origin bottom-left, y up]."""
    import shlex
    txt = subprocess.run(["dot", "-Tplain", str(dot_path)],
                         capture_output=True, text=True, check=True).stdout
    gw = gh = 1.0
    nodes, edges = {}, []
    for line in txt.splitlines():
        p = shlex.split(line)
        if not p:
            continue
        if p[0] == "graph":
            gw, gh = float(p[2]), float(p[3])
        elif p[0] == "node":
            nodes[p[1]] = (float(p[2]), float(p[3]), float(p[4]), float(p[5]))
        elif p[0] == "edge":
            edges.append((p[1], p[2]))
    return gw, gh, nodes, edges


def to_pptx(dot_path, pptx_path):
    """Rebuild the dot layout as NATIVE, editable PowerPoint shapes (ovals =
    nodes, rounded rects = topics) with connectors bound to the shapes."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn

    gw, gh, nodes, edges = _parse_plain(dot_path)
    topics = {t for n in NODES.values() for t in (n["pub"] + n["sub"])}

    MX, MTOP, MBOT = 0.5, 1.0, 0.4
    prs = Presentation()
    prs.slide_width = Inches(gw + 2 * MX)
    prs.slide_height = Inches(gh + MTOP + MBOT)
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

    def rgb(hexs):
        return RGBColor.from_string(hexs.lstrip("#"))

    def to_xy(cx, cy, w, h):                 # dot center -> pptx top-left (inches)
        return (MX + cx - w / 2, MTOP + (gh - cy - h / 2))

    # ── Gazebo/ros2_control cluster box (bbox of its member nodes) ──────────
    gz = [nodes[n] for n in GAZEBO_CLUSTER if n in nodes]
    if gz:
        x0 = min(cx - w / 2 for cx, cy, w, h in gz) - 0.18
        x1 = max(cx + w / 2 for cx, cy, w, h in gz) + 0.18
        y0 = min(cy - h / 2 for cx, cy, w, h in gz) - 0.18
        y1 = max(cy + h / 2 for cx, cy, w, h in gz) + 0.40   # room for label
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(MX + x0),
            Inches(MTOP + (gh - y1)), Inches(x1 - x0), Inches(y1 - y0))
        box.fill.solid(); box.fill.fore_color.rgb = rgb("#eef7ee")
        box.line.color.rgb = rgb("#2e7d32"); box.line.width = Pt(1.5)
        box.shadow.inherit = False
        tf = box.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.LEFT
        r = p0.add_run(); r.text = "Gazebo / ros2_control"
        r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = rgb("#2e7d32")
        p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run(); r1.text = "controller_manager"
        r1.font.size = Pt(8); r1.font.italic = True; r1.font.color.rgb = rgb("#4a7a4a")

    # ── nodes (ovals) + topics (rounded rects) ──────────────────────────────
    shp = {}
    for name, (cx, cy, w, h) in nodes.items():
        left, top = to_xy(cx, cy, w, h)
        is_node = name in NODES
        kind = MSO_SHAPE.OVAL if is_node else MSO_SHAPE.ROUNDED_RECTANGLE
        s = slide.shapes.add_shape(kind, Inches(left), Inches(top),
                                   Inches(w), Inches(h))
        if name in GAZEBO_CLUSTER:
            fc, lc = "#cdeccd", "#2e7d32"
        elif is_node:
            fc, lc = "#cfe2ff", "#3a6ea5"
        else:
            fc, lc = "#fff3c4", "#b58900"
        s.fill.solid(); s.fill.fore_color.rgb = rgb(fc)
        s.line.color.rgb = rgb(lc); s.line.width = Pt(1.25)
        s.shadow.inherit = False
        tf = s.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(2)
        tf.margin_top = tf.margin_bottom = Pt(1)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = name
        run.font.size = Pt(9); run.font.color.rgb = rgb("#1b1b1b")
        shp[name] = s

    # ── connectors (bound to shapes so they follow when dragged) ────────────
    # cardinal connection-site indices for an autoshape: 0=top 1=left 2=bottom 3=right
    def side(frm, to):
        fx, fy, _, _ = nodes[frm]; tx, ty, _, _ = nodes[to]
        dx, dy = tx - fx, ty - fy
        if abs(dx) >= abs(dy):
            return 3 if dx >= 0 else 1
        return 0 if dy >= 0 else 2     # dot y is up: dy>0 -> 'to' is above -> top
    for a, b in edges:
        if a not in shp or b not in shp:
            continue
        cn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(MX), Inches(MTOP),
            Inches(MX + 1), Inches(MTOP + 1))
        cn.line.color.rgb = rgb("#5b6470"); cn.line.width = Pt(1.1)
        cn.begin_connect(shp[a], side(a, b))
        cn.end_connect(shp[b], side(b, a))
        ln = cn.line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:tailEnd"),
                                 {"type": "triangle", "w": "med", "len": "med"}))

    # ── title ───────────────────────────────────────────────────────────────
    tb = slide.shapes.add_textbox(Inches(MX), Inches(0.18),
                                  Inches(gw + MX), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "ROS 2 computation graph — pendant backend (Gazebo mode)"
    r.font.bold = True; r.font.size = Pt(18)
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = "rqt_graph Nodes/Topics view  ·  editable"
    r2.font.size = Pt(11); r2.font.color.rgb = rgb("#555555")

    prs.save(str(pptx_path))


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="system_rqt_graph",
                    help="output prefix (writes .dot/.png/.pdf and .pptx)")
    ap.add_argument("--no-pptx", action="store_true",
                    help="skip the editable PowerPoint output")
    args = ap.parse_args()
    out = Path(args.out)
    out.parent.mkdir(parents=True, exist_ok=True)
    dot = build_dot()
    out.with_suffix(".dot").write_text(dot)
    for fmt in ("png", "pdf"):
        subprocess.run(["dot", "-T" + fmt, "-Gdpi=150",
                        str(out.with_suffix(".dot")),
                        "-o", str(out.with_suffix("." + fmt))], check=True)
    msg = f"wrote {out}.dot/.png/.pdf"
    if not args.no_pptx:
        to_pptx(out.with_suffix(".dot"), out.with_suffix(".pptx"))
        msg += "/.pptx"
    print(f"{msg}  ({len(NODES)} nodes, "
          f"{len({t for n in NODES.values() for t in n['pub']+n['sub']})} topics)")


if __name__ == "__main__":
    main()
